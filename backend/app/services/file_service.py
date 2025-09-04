import os
import uuid
from typing import Dict, List
from fastapi import UploadFile
from langchain_community.vectorstores import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA, LLMChain
from langchain_openai.chat_models.base import BaseChatOpenAI
from langchain.docstore.document import Document
from queue import Queue
import threading
import pdfplumber
import pandas as pd
from docx import Document
from pptx import Presentation
from unstructured.partition.auto import partition
import jieba
from collections import Counter
import re
from datetime import datetime
from app.services.db_service import DatabaseService
from app.database import SessionLocal
from app.utils.llm_helper import llm_helper 

class FileService:
    def __init__(self):
        self.upload_dir = "uploads"
        self.supported_formats = {'.pdf', '.docx', '.xlsx', '.xls', '.txt', '.pptx'}
        # 确保上传目录存在
        os.makedirs(self.upload_dir, exist_ok=True)
        self.db_service = None
        # 🔥 持久化目录
        self.vector_dir = "./vectorstores"
        self.embeddings = HuggingFaceEmbeddings(model_name="shibing624/text2vec-base-chinese")
        
    def _get_db_service(self):
        """获取数据库服务实例"""
        if not self.db_service:
            db = SessionLocal()
            self.db_service = DatabaseService(db)
        return self.db_service
    # =========================
    # 🔥 1. 构建向量库
    # =========================
    def build_vector_store(self, file_id: str, full_text: str):
        """为某个文件构建向量库"""
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        docs = [Document(page_content=chunk, metadata={"file_id": file_id})
                for chunk in splitter.split_text(full_text)]
        
        vector_db = Chroma.from_documents(
            documents=docs,
            embedding=self.embeddings,
            persist_directory=os.path.join(self.vector_dir, file_id)
        )
        vector_db.persist()
        return True
    
    def _get_vectorstore(self, file_id: str):
        """加载文件的向量库"""
        return Chroma(
            embedding_function=self.embeddings,
            persist_directory=os.path.join(self.vector_dir, file_id)
        )
    # =========================
    # 🔥 2. 非流式文件问答
    # =========================
    def ask_file(self, file_id: str, question: str) -> dict:
        vector_db = self._get_vectorstore(file_id)
        retriever = vector_db.as_retriever(search_type="mmr", search_kwargs={"k": 5})

        

        PROMPT_TEMPLATE = """
        你是一个专业的知识助手，请严格根据提供的上下文信息回答问题。
        如果上下文不包含答案，请回答 "根据现有知识无法回答该问题"。

        上下文：
        {context}

        问题：{question}
        """

        qa_chain = RetrievalQA.from_chain_type(
            llm=llm_helper.llm,
            chain_type="stuff",
            retriever=retriever,
            chain_type_kwargs={"prompt": PromptTemplate(
                template=PROMPT_TEMPLATE,
                input_variables=["context", "question"]
            )},
            return_source_documents=True
        )

        result = qa_chain({"query": question})
        return {
            "answer": result["result"],
            "sources": [
                {
                    "page_content": doc.page_content[:200],
                    "metadata": doc.metadata
                } for doc in result.get("source_documents", [])
            ]
        }
    # =========================
    # 🔥 3. 流式问答（带对话历史）
    # =========================
    def ask_file_stream(self, file_id: str, messages: List[dict], summary_text: str):
        vector_db = self._get_vectorstore(file_id)
        retriever = vector_db.as_retriever(search_type="mmr", search_kwargs={"k": 5})
        
        query = messages[-1]["content"]
        docs = retriever.get_relevant_documents(query)
        context = "\n".join([doc.page_content for doc in docs])
        history = "\n".join([f"{m['role']}：{m['content']}" for m in messages[:-1]])

        PROMPT_TEMPLATE = """
        你是一个专业的知识问答助手，请根据以下知识上下文和对话历史回答用户问题。
        如果无法根据内容回答，请回复 "根据现有知识无法回答该问题"。

        【知识上下文】
        {context}

        【历史摘要】
        {summary}

        【历史对话】
        {history}

        【当前问题】
        {question}
        """
        chain = LLMChain(
            llm=llm_helper.llm,
            prompt=PromptTemplate(
                input_variables=["context", "summary", "history", "question"],
                template=PROMPT_TEMPLATE
            )
        )

        q = Queue()

        def _callback(token):
            q.put(token)

        def run_chain():
            try:
                chain.run({
                    "context": context,
                    "summary": summary_text,
                    "history": history,
                    "question": query
                })
            finally:
                q.put(None)

        threading.Thread(target=run_chain).start()

        def token_stream():
            while True:
                token = q.get()
                if token is None:
                    break
                yield f"data: {token}\n\n"

        return token_stream()

    async def save_upload_file(self, file: UploadFile, conversation_id: str) -> str:
        """保存上传的文件"""
        try:
            file_extension = os.path.splitext(file.filename)[1].lower()
            unique_filename = f"{conversation_id}_{uuid.uuid4()}{file_extension}"
            file_path = os.path.join(self.upload_dir, unique_filename)
            
            with open(file_path, "wb") as buffer:
                content = await file.read()
                buffer.write(content)
            
            return file_path
        except Exception as e:
            raise Exception(f"保存文件失败: {str(e)}")
    
    def get_file_info(self, file_path: str, original_name: str) -> Dict:
        """获取文件基本信息"""
        try:
            file_stats = os.stat(file_path)
            file_size = file_stats.st_size
            file_extension = os.path.splitext(file_path)[1].lower()
            
            file_info = {
                "file_name": original_name,
                "file_path": file_path,
                "file_size": file_size,
                "file_format": file_extension,
                "upload_time": datetime.now().isoformat()
            }
            
            # 根据文件类型获取特定信息
            if file_extension == '.pdf':
                info = self._get_pdf_info(file_path)
                file_info.update(info)
            elif file_extension == '.docx':
                info = self._get_docx_info(file_path)
                file_info.update(info)
            elif file_extension in ['.xlsx', '.xls']:
                info = self._get_excel_info(file_path)
                file_info.update(info)
            elif file_extension == '.pptx':
                info = self._get_pptx_info(file_path)
                file_info.update(info)
            elif file_extension == '.txt':
                info = self._get_txt_info(file_path)
                file_info.update(info)
            
            return file_info
        except Exception as e:
            raise Exception(f"获取文件信息失败: {str(e)}")
    
    def _get_pdf_info(self, file_path: str) -> Dict:
        try:
            with pdfplumber.open(file_path) as pdf:
                page_count = len(pdf.pages)
                total_text = "".join([page.extract_text() or "" for page in pdf.pages])
                return {
                    "page_count": page_count,
                    "word_count": len(total_text.split()),
                    "character_count": len(total_text)
                }
        except Exception as e:
            return {"error": f"读取PDF失败: {str(e)}"}
    
    def _get_docx_info(self, file_path: str) -> Dict:
        try:
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            total_text = "\n".join(paragraphs)
            return {
                "word_count": len(total_text.split()),
                "character_count": len(total_text),
                "title": doc.core_properties.title or "",
                "author": doc.core_properties.author or ""
            }
        except Exception as e:
            return {"error": f"读取DOCX失败: {str(e)}"}
    
    def _get_excel_info(self, file_path: str) -> Dict:
        try:
            excel_file = pd.ExcelFile(file_path)
            sheet_names = excel_file.sheet_names
            total_rows, total_cols = 0, 0
            for sheet_name in sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                total_rows += len(df)
                total_cols = max(total_cols, len(df.columns))
            return {
                "sheet_count": len(sheet_names),
                "sheet_names": sheet_names,
                "total_rows": total_rows,
                "total_columns": total_cols
            }
        except Exception as e:
            return {"error": f"读取Excel失败: {str(e)}"}
    
    def _get_pptx_info(self, file_path: str) -> Dict:
        try:
            prs = Presentation(file_path)
            slide_count = len(prs.slides)
            total_text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        total_text += shape.text + " "
            return {
                "slide_count": slide_count,
                "word_count": len(total_text.split()),
                "character_count": len(total_text)
            }
        except Exception as e:
            return {"error": f"读取PPTX失败: {str(e)}"}
    
    def _get_txt_info(self, file_path: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return {
                "word_count": len(content.split()),
                "character_count": len(content)
            }
        except Exception as e:
            return {"error": f"读取TXT失败: {str(e)}"}
    
    def extract_content(self, file_path: str) -> Dict:
        """提取文件内容"""
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension == '.pdf':
                return self._extract_pdf_content(file_path)
            elif file_extension == '.docx':
                return self._extract_docx_content(file_path)
            elif file_extension in ['.xlsx', '.xls']:
                return self._extract_excel_content(file_path)
            elif file_extension == '.pptx':
                return self._extract_pptx_content(file_path)
            elif file_extension == '.txt':
                return self._extract_txt_content(file_path)
            else:
                # 改为使用 unstructured 兜底
                return self._extract_with_unstructured(file_path)
        except Exception as e:
            raise Exception(f"提取文件内容失败: {str(e)}")
    
    def _extract_pdf_content(self, file_path: str) -> Dict:
        try:
            content_data = {"pages": [], "full_text": "", "tables": [], "metadata": {}}
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text() or ""
                    content_data["pages"].append({
                        "page_number": i + 1,
                        "content": page_text
                    })
                    content_data["full_text"] += page_text + "\n"
                # 表格提取
                for i, page in enumerate(pdf.pages):
                    tables = page.extract_tables()
                    if tables:
                        content_data["tables"].extend([{"page": i + 1, "table": t} for t in tables])
            return content_data
        except Exception as e:
            return {"error": f"提取PDF内容失败: {str(e)}"}
    
    def _extract_docx_content(self, file_path: str) -> Dict:
        try:
            doc = Document(file_path)
            paragraphs, full_text = [], ""
            for i, p in enumerate(doc.paragraphs):
                if p.text.strip():
                    paragraphs.append({
                        "paragraph_number": i + 1,
                        "content": p.text,
                        "style": p.style.name if p.style else ""
                    })
                    full_text += p.text + "\n"
            tables = []
            for i, t in enumerate(doc.tables):
                data = [[c.text for c in r.cells] for r in t.rows]
                tables.append({"table_number": i + 1, "data": data})
            return {
                "paragraphs": paragraphs,
                "full_text": full_text,
                "tables": tables,
                "metadata": {
                    "title": doc.core_properties.title or "",
                    "author": doc.core_properties.author or "",
                }
            }
        except Exception as e:
            return {"error": f"提取DOCX内容失败: {str(e)}"}
    
    def _extract_excel_content(self, file_path: str) -> Dict:
        try:
            excel_file = pd.ExcelFile(file_path)
            sheets_data = {}
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                sheets_data[sheet_name] = {
                    "data": df.to_dict('records'),
                    "columns": df.columns.tolist(),
                    "row_count": len(df)
                }
            return {"sheets": sheets_data, "sheet_names": excel_file.sheet_names}
        except Exception as e:
            return {"error": f"提取Excel内容失败: {str(e)}"}
    
    def _extract_pptx_content(self, file_path: str) -> Dict:
        try:
            prs = Presentation(file_path)
            slides_data, full_text = [], ""
            for i, slide in enumerate(prs.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + " "
                slides_data.append({"slide_number": i + 1, "content": slide_text.strip()})
                full_text += slide_text + "\n"
            return {"slides": slides_data, "full_text": full_text, "slide_count": len(slides_data)}
        except Exception as e:
            return {"error": f"提取PPTX内容失败: {str(e)}"}
    
    def _extract_txt_content(self, file_path: str) -> Dict:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            lines = content.split('\n')
            paragraphs = [{"paragraph_number": i + 1, "content": line} for i, line in enumerate(lines) if line.strip()]
            return {"paragraphs": paragraphs, "full_text": content}
        except Exception as e:
            return {"error": f"提取TXT内容失败: {str(e)}"}
    
    def _extract_with_unstructured(self, file_path: str) -> Dict:
        """使用 unstructured 兜底提取"""
        try:
            elements = partition(filename=file_path)
            texts = [el.text for el in elements if hasattr(el, "text") and el.text]
            return {
                "full_text": "\n".join(texts),
                "paragraphs": [{"content": t} for t in texts]
            }
        except Exception as e:
            return {"error": f"使用unstructured提取内容失败: {str(e)}"}
    async def _generate_llm_summary(self, text: str,max_length: int = 300) -> str:
        """
        使用 DeepSeek 生成自然摘要（通过 llm_helper）
        """
        if not text.strip():
            return "文本为空，无法生成摘要"

        # 截断，避免输入过长
        snippet = text[:2000]

        messages = [
            {"role": "system", "content": "你是一名专业的文档分析助手，请帮我生成简明扼要的摘要。"},
            {"role": "user", "content": f"请为以下文本生成一个不超过{max_length}字的摘要：\n\n{snippet}"}
        ]

        try:
            response = await llm_helper.chat_completion(messages)
            summary = response["choices"][0]["message"]["content"].strip()
            return summary
        except Exception as e:
            return f"LLM摘要生成失败: {str(e)}"

    async def analyze_content(self, content_data: Dict, file_info: Dict) -> Dict:
        try:
            full_text = content_data.get("full_text", "")
            word_count, character_count = len(full_text.split()), len(full_text)
            keywords = self._extract_keywords(full_text)
            entities = self._extract_entities(full_text)
            summaries = {
                "summary_100": self._generate_summary(full_text, 100),
                "summary_300": self._generate_summary(full_text, 300),
                "summary_1000": self._generate_summary(full_text, 1000),
            }
            # 可选：调用大模型做更自然的摘要（如果接了 LLM）
            try:
                llm_summary = await self._generate_llm_summary(full_text)
                summaries["llm_summary"] = llm_summary
            except Exception:
                summaries["llm_summary"] = "LLM摘要不可用"

            return {
                "statistics": {
                    "word_count": word_count,
                    "character_count": character_count,
                    "paragraph_count": len(content_data.get("paragraphs", [])),
                    "page_count": file_info.get("page_count", 0),
                    "slide_count": file_info.get("slide_count", 0)
                },
                "keywords": keywords,
                "entities": entities,
                "summaries": summaries
            }
        except Exception as e:
            raise Exception(f"分析文件内容失败: {str(e)}")
    
    def _extract_keywords(self, text: str) -> List[Dict]:
        try:
            words = jieba.lcut(text)
            filtered = [w for w in words if len(w) > 1 and w.isalpha()]
            freq = Counter(filtered)
            return [{"word": w, "frequency": f} for w, f in freq.most_common(20)]
        except Exception:
            return [{"word": "关键词提取失败", "frequency": 0}]
    
    def _extract_entities(self, text: str) -> Dict:
        try:
            return {
                "dates": re.findall(r'\d{4}[-年]\d{1,2}[-月]\d{1,2}[日]?', text),
                "numbers": re.findall(r'\d+(?:\.\d+)?', text),
                "organizations": [],
                "persons": []
            }
        except Exception as e:
            return {"error": f"实体识别失败: {str(e)}"}
    
    def _generate_summary(self, text: str, max_length: int) -> str:
        try:
            sentences = re.split(r'[。！？.!?]', text)
            summary, length = "", 0
            for s in sentences:
                if s.strip() and length + len(s) <= max_length:
                    summary += s + "。"
                    length += len(s)
                else:
                    break
            return summary if summary else (text[:max_length] + "..." if len(text) > max_length else text)
        except Exception:
            return f"摘要生成失败"
    
    async def process_upload_and_analyze(self, file: UploadFile, conversation_id: str) -> Dict:
        try:
            file_path = await self.save_upload_file(file, conversation_id)
            file_info = self.get_file_info(file_path, file.filename)
            content_data = self.extract_content(file_path)
            analysis_data = await self.analyze_content(content_data, file_info)
            insights = self._generate_insights(file_info, analysis_data, content_data)
            
            # 🔥 自动构建向量库
            if "full_text" in content_data and content_data["full_text"].strip():
                file_id = os.path.splitext(os.path.basename(file_path))[0]
                self.build_vector_store(file_id, content_data["full_text"])
            
            return {
                "file_path": file_path,
                "file_info": file_info,
                "content_data": content_data,
                "analysis_data": analysis_data,
                "insights": insights,
                "success": True
            }
        except Exception as e:
            return {"error": str(e), "success": False}
    
    def _generate_insights(self, file_info: Dict, analysis_data: Dict, content_data: Dict) -> str:
        try:
            insights = f"""📄 文件分析洞察：

📊 基本信息：
- 文件名: {file_info.get('file_name', '未知')}
- 文件大小: {file_info.get('file_size', 0)} 字节
- 文件格式: {file_info.get('file_format', '未知')}
- 上传时间: {file_info.get('upload_time', '未知')}

📈 统计信息：
- 字数: {analysis_data.get('statistics', {}).get('word_count', 0)}
- 字符数: {analysis_data.get('statistics', {}).get('character_count', 0)}
"""
            if file_info.get('file_format') == '.pdf':
                insights += f"- 页数: {file_info.get('page_count', 0)}\n"
            elif file_info.get('file_format') == '.pptx':
                insights += f"- 幻灯片数: {file_info.get('slide_count', 0)}\n"
            elif file_info.get('file_format') in ['.xlsx', '.xls']:
                insights += f"- 工作表数: {file_info.get('sheet_count', 0)}\n"
            keywords = analysis_data.get('keywords', [])
            if keywords:
                insights += "\n🔑 关键词：\n"
                for i, kw in enumerate(keywords[:10]):
                    insights += f"  {i+1}. {kw['word']} ({kw['frequency']}次)\n"
            summaries = analysis_data.get('summaries', {})
            if summaries.get('summary_100'):
                insights += f"\n📋 100字摘要：\n{summaries['summary_100']}\n"
            if summaries.get('llm_summary'):
                insights += f"\n🤖 LLM摘要：\n{summaries['llm_summary']}\n"
            return insights
        except Exception:
            return "生成洞察失败"
# 全局文件服务实例
file_service = FileService()